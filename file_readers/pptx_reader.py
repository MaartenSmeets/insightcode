# file_readers/pptx_reader.py
import logging
from pptx import Presentation

FILE_EXTENSIONS = ['.pptx']

def read_file(file_path):
    """Read contents from a .pptx file."""
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"Error reading pptx file {file_path}: {e}")
        return ""